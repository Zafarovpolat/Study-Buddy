# backend/app/services/presentation_service.py
import json
import re
from typing import Dict, Any, List, Optional
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor  # Правильный импорт!

from app.services.ai_service import gemini_service


class PresentationService:
    """Сервис генерации презентаций"""
    
    # Цветовые схемы
    THEMES = {
        "blue": {
            "primary": RGBColor(0x1E, 0x40, 0xAF),
            "secondary": RGBColor(0x3B, 0x82, 0xF6),
            "accent": RGBColor(0x60, 0xA5, 0xFA),
            "text": RGBColor(0x1F, 0x29, 0x37),
            "light": RGBColor(0xEF, 0xF6, 0xFF),
        },
        "green": {
            "primary": RGBColor(0x05, 0x96, 0x69),
            "secondary": RGBColor(0x10, 0xB9, 0x81),
            "accent": RGBColor(0x34, 0xD3, 0x99),
            "text": RGBColor(0x1F, 0x29, 0x37),
            "light": RGBColor(0xEC, 0xFD, 0xF5),
        },
        "purple": {
            "primary": RGBColor(0x7C, 0x3A, 0xED),
            "secondary": RGBColor(0x8B, 0x5C, 0xF6),
            "accent": RGBColor(0xA7, 0x8B, 0xFA),
            "text": RGBColor(0x1F, 0x29, 0x37),
            "light": RGBColor(0xF5, 0xF3, 0xFF),
        },
        "orange": {
            "primary": RGBColor(0xEA, 0x58, 0x0C),
            "secondary": RGBColor(0xF9, 0x73, 0x16),
            "accent": RGBColor(0xFB, 0x92, 0x3C),
            "text": RGBColor(0x1F, 0x29, 0x37),
            "light": RGBColor(0xFF, 0xF7, 0xED),
        },
    }
    
    async def generate_presentation_structure(
        self, 
        topic: str, 
        num_slides: int = 10,
        style: str = "professional"
    ) -> Dict[str, Any]:
        """Генерирует структуру презентации через AI"""
        
        style_instructions = {
            "professional": "деловой и информативный стиль, факты и данные",
            "educational": "образовательный стиль, простые объяснения, примеры",
            "creative": "креативный стиль, метафоры, интересные факты",
            "minimal": "минималистичный стиль, только ключевые тезисы",
        }
        
        prompt = f"""Создай структуру презентации на тему: "{topic}"

Количество слайдов: {num_slides}
Стиль: {style_instructions.get(style, style_instructions["professional"])}

Формат JSON:
{{
  "title": "Название презентации",
  "subtitle": "Подзаголовок",
  "author": "Lecto AI",
  "slides": [
    {{
      "type": "title",
      "title": "Название презентации",
      "subtitle": "Подзаголовок"
    }},
    {{
      "type": "content",
      "title": "Заголовок слайда",
      "bullets": ["Пункт 1", "Пункт 2", "Пункт 3"],
      "notes": "Заметки докладчика (подробнее)"
    }},
    {{
      "type": "two_columns",
      "title": "Сравнение",
      "left_title": "Левая колонка",
      "left_bullets": ["Пункт 1", "Пункт 2"],
      "right_title": "Правая колонка", 
      "right_bullets": ["Пункт 1", "Пункт 2"],
      "notes": "Заметки"
    }},
    {{
      "type": "quote",
      "quote": "Цитата или ключевая мысль",
      "author": "Автор цитаты",
      "notes": "Контекст цитаты"
    }},
    {{
      "type": "conclusion",
      "title": "Заключение",
      "bullets": ["Вывод 1", "Вывод 2", "Вывод 3"],
      "call_to_action": "Призыв к действию"
    }}
  ]
}}

Требования:
1. Первый слайд - титульный (type: "title")
2. Последний слайд - заключение (type: "conclusion")
3. Используй разные типы слайдов для разнообразия
4. Каждый слайд: 3-5 пунктов максимум
5. Заметки докладчика для каждого слайда
6. Логичная структура: введение - основная часть - заключение

Верни ТОЛЬКО валидный JSON!"""

        try:
            text = await gemini_service._generate_async(prompt)
            text = text.strip()
            
            text = re.sub(r'^```json\s*', '', text)
            text = re.sub(r'^```\s*', '', text)
            text = re.sub(r'\s*```$', '', text)
            
            structure = json.loads(text)
            
            if "slides" not in structure or not structure["slides"]:
                raise ValueError("No slides generated")
            
            return structure
            
        except json.JSONDecodeError as e:
            print(f"Presentation JSON error: {e}")
            return self._get_fallback_structure(topic)
        except Exception as e:
            print(f"Presentation generation error: {e}")
            raise
    
    def _get_fallback_structure(self, topic: str) -> Dict[str, Any]:
        """Запасная структура при ошибке AI"""
        return {
            "title": topic,
            "subtitle": "Презентация",
            "author": "Lecto AI",
            "slides": [
                {
                    "type": "title",
                    "title": topic,
                    "subtitle": "Презентация создана с помощью AI"
                },
                {
                    "type": "content",
                    "title": "Введение",
                    "bullets": [
                        "Основные понятия темы",
                        "Почему это важно",
                        "Что мы рассмотрим"
                    ],
                    "notes": "Начните с общего обзора темы"
                },
                {
                    "type": "conclusion",
                    "title": "Заключение",
                    "bullets": [
                        "Ключевые выводы",
                        "Практическое применение"
                    ],
                    "call_to_action": "Изучите тему подробнее!"
                }
            ]
        }
    
    def create_pptx(
        self, 
        structure: Dict[str, Any],
        theme: str = "blue"
    ) -> BytesIO:
        """Создает PPTX файл из структуры"""
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        colors = self.THEMES.get(theme, self.THEMES["blue"])
        
        for slide_data in structure.get("slides", []):
            slide_type = slide_data.get("type", "content")
            
            if slide_type == "title":
                self._add_title_slide(prs, slide_data, colors)
            elif slide_type == "two_columns":
                self._add_two_column_slide(prs, slide_data, colors)
            elif slide_type == "quote":
                self._add_quote_slide(prs, slide_data, colors)
            elif slide_type == "conclusion":
                self._add_conclusion_slide(prs, slide_data, colors)
            else:
                self._add_content_slide(prs, slide_data, colors)
        
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output
    
    def _add_title_slide(self, prs, data: Dict, colors: Dict):
        """Титульный слайд"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2.5),
            Inches(13.333), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["primary"]
        shape.line.fill.background()
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.7),
            Inches(12.333), Inches(1.2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "Презентация")
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        if data.get("subtitle"):
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4),
                Inches(12.333), Inches(0.8)
            )
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = data.get("subtitle", "")
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER
    
    def _add_content_slide(self, prs, data: Dict, colors: Dict):
        """Слайд с контентом"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = colors["primary"]
        header_shape.line.fill.background()
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.6),
            Inches(12), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(data.get("bullets", [])):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "* " + bullet
            p.font.size = Pt(24)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(12)
            p.space_after = Pt(6)
        
        if data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = data["notes"]
    
    def _add_two_column_slide(self, prs, data: Dict, colors: Dict):
        """Слайд с двумя колонками"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = colors["primary"]
        header_shape.line.fill.background()
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        left_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(6), Inches(0.6)
        )
        tf = left_title.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("left_title", "")
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = colors["secondary"]
        
        left_content = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.2),
            Inches(6), Inches(4.5)
        )
        tf = left_content.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(data.get("left_bullets", [])):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "* " + bullet
            p.font.size = Pt(20)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(8)
        
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.5), Inches(1.5),
            Inches(0.05), Inches(5.5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = colors["accent"]
        divider.line.fill.background()
        
        right_title = slide.shapes.add_textbox(
            Inches(6.833), Inches(1.5),
            Inches(6), Inches(0.6)
        )
        tf = right_title.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("right_title", "")
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = colors["secondary"]
        
        right_content = slide.shapes.add_textbox(
            Inches(6.833), Inches(2.2),
            Inches(6), Inches(4.5)
        )
        tf = right_content.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(data.get("right_bullets", [])):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "* " + bullet
            p.font.size = Pt(20)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(8)
        
        if data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = data["notes"]
    
    def _add_quote_slide(self, prs, data: Dict, colors: Dict):
        """Слайд с цитатой"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = colors["light"]
        bg_shape.line.fill.background()
        
        quote_mark = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(2), Inches(2)
        )
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = '"'
        p.font.size = Pt(120)
        p.font.color.rgb = colors["accent"]
        
        quote_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5),
            Inches(10.333), Inches(3)
        )
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("quote", "")
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.color.rgb = colors["text"]
        p.alignment = PP_ALIGN.CENTER
        
        if data.get("author"):
            author_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(5.5),
                Inches(10.333), Inches(0.8)
            )
            tf = author_box.text_frame
            p = tf.paragraphs[0]
            p.text = "-- " + data["author"]
            p.font.size = Pt(20)
            p.font.color.rgb = colors["secondary"]
            p.alignment = PP_ALIGN.RIGHT
        
        if data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = data["notes"]
    
    def _add_conclusion_slide(self, prs, data: Dict, colors: Dict):
        """Заключительный слайд"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = colors["primary"]
        header_shape.line.fill.background()
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "Заключение")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.6),
            Inches(12), Inches(4)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(data.get("bullets", [])):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "[OK] " + bullet
            p.font.size = Pt(24)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(12)
        
        if data.get("call_to_action"):
            cta_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6),
                Inches(12.333), Inches(1)
            )
            tf = cta_box.text_frame
            p = tf.paragraphs[0]
            p.text = data["call_to_action"]
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = colors["secondary"]
            p.alignment = PP_ALIGN.CENTER


presentation_service = PresentationService()